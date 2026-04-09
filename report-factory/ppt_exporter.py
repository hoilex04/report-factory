import sys
import json
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# --- 老板的 Obsidian 路径配置 ---
OBSIDIAN_VAULT = r"D:\003_Resource\04_Obsidian\Atomic-card"
CARDS_DIR = os.path.join(OBSIDIAN_VAULT, "Cards")
CANVAS_FILE = os.path.join(OBSIDIAN_VAULT, "Inbox_Dashboard.canvas")

def ensure_obsidian_dirs():
    """确保 Obsidian 的存储目录存在"""
    if not os.path.exists(CARDS_DIR):
        os.makedirs(CARDS_DIR)

def update_master_index(new_blocks, index_path='master_index.json'):
    """维护全局索引文件"""
    index_data = []
    if os.path.exists(index_path):
        with open(index_path, 'r', encoding='utf-8') as f:
            index_data = json.load(f)
    index_data.extend(new_blocks)
    with open(index_path, 'w', encoding='utf-8') as f:
        json.dump(index_data, f, ensure_ascii=False, indent=2)

def extract_date_from_uid(uid):
    """从 UID 中提取日期，格式：AIH-20260208-01 -> 20260208"""
    parts = uid.split('-')
    if len(parts) >= 2:
        return parts[1]  # 返回 20260208
    return datetime.now().strftime('%Y%m%d')

def generate_new_filename(item):
    """生成新的文件名：{domain}_{title}_{YYYYMMDD}.md"""
    domain = item.get('domain', 'Unknown')
    # 简化 domain 名称
    if domain == 'AI-Hardware':
        domain = 'AIH'
    elif domain == 'Embodied-AI':
        domain = 'EAI'

    title = item.get('title', 'Untitled')
    uid = item.get('uid', '')
    date = extract_date_from_uid(uid)

    # 清理标题中的特殊字符（Windows 文件名不允许的字符）
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in invalid_chars:
        title = title.replace(char, '')

    filename = f"{domain}_{title}_{date}.md"
    return filename

def sync_to_obsidian(item):
    """同步到 Obsidian：创建 MD 文件（新命名规则）"""
    uid = item.get('uid')
    title = item.get('title', 'Untitled')
    domain = item.get('domain', 'General')

    # 构建 quantitative_data 部分（如果存在）
    quant_section = ""
    if 'quantitative_data' in item:
        quant_section = "\n\n## 📊 定量数据\n\n```json\n" + json.dumps(item['quantitative_data'], ensure_ascii=False, indent=2) + "\n```\n"

    # 使用新的文件名生成规则
    md_filename = generate_new_filename(item)
    md_path = os.path.join(CARDS_DIR, md_filename)

    # 优化 content 排版：将长文本按句号、分号分段
    content = item.get('content', '')
    # 在中文句号、分号、冒号后添加换行（但保留在Callout引用符号内）
    formatted_content = content.replace('。', '。\n> \n> ')
    formatted_content = formatted_content.replace('；', '；\n> \n> ')
    # 移除末尾多余的引用符号
    formatted_content = formatted_content.rstrip('\n> ')

    # 写入内容
    md_content = f"""---
uid: {uid}
domain: {domain}
tags: {item.get('tags')}
source: {item.get('source_short')}
date: {datetime.now().strftime('%Y-%m-%d')}
---
# {title}

> [!abstract] 核心论据
> {formatted_content}
{quant_section}
**视觉建议：** {item.get('visual_suggestion')}

**原始来源：** {item.get('source')}
"""
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write(md_content)

    # 更新 Canvas (使用新的文件名)
    update_canvas(uid, md_filename)
    
def update_canvas(uid, md_filename):
    """在 Obsidian Canvas 中自动添加新卡片节点"""
    canvas_data = {"nodes": [], "edges": []}
    if os.path.exists(CANVAS_FILE):
        try:
            with open(CANVAS_FILE, 'r', encoding='utf-8') as f:
                canvas_data = json.load(f)
        except:
            pass
    
    # 简单的堆叠布局：新节点向下排列
    node_spacing = 500
    current_nodes = len(canvas_data.get('nodes', []))
    
    new_node = {
        "id": uid,
        "type": "file",
        "file": f"Cards/{md_filename}",
        "x": 0,
        "y": current_nodes * node_spacing,
        "width": 400,
        "height": 450
    }
    
    if "nodes" not in canvas_data: canvas_data["nodes"] = []
    canvas_data["nodes"].append(new_node)
    
    with open(CANVAS_FILE, 'w', encoding='utf-8') as f:
        json.dump(canvas_data, f, ensure_ascii=False, indent=2)

def create_ppt(data_json_path):
    # 读取数据
    with open(data_json_path, 'r', encoding='utf-8') as f:
        blocks = json.load(f)

    # 1. 更新全局索引
    update_master_index(blocks)
    
    # 2. 准备 Obsidian 环境
    ensure_obsidian_dirs()

    # 3. 生成 PPT 并 同步 Obsidian
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    for item in blocks:
        # --- 同步 Obsidian ---
        sync_to_obsidian(item)

        # --- 生成 PPT 豆腐块 (保留老板原有版式) ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 绘制卡片背景
        shape = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(252, 252, 252)
        shape.line.color.rgb = RGBColor(200, 200, 200)

        # 标题
        t = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        p = t.text_frame.paragraphs[0]
        p.text = f"[{item.get('uid')}] {item.get('title')}"
        p.font.bold, p.font.size, p.font.color.rgb = True, Pt(26), RGBColor(0, 51, 102)

        # 正文
        b = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.5), Inches(3.5))
        tf_b = b.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = item.get('content', '')
        p_b.font.size, p_b.font.color.rgb = Pt(18), RGBColor(60, 60, 60)

        # 右侧视觉建议
        v = slide.shapes.add_textbox(Inches(7.8), Inches(2.0), Inches(4.5), Inches(3.5))
        p_v = v.text_frame.paragraphs[0]
        p_v.text = f"【配图建议】\n{item.get('visual_suggestion')}"
        p_i = v.text_frame.paragraphs[0].font
        p_i.size, p_i.italic, p_i.color.rgb = Pt(14), True, RGBColor(150, 150, 150)

        # 底部元数据栏
        m = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5))
        p_m = m.text_frame.paragraphs[0]
        p_m.text = f"Tags: {item.get('tags')}  |  Source: {item.get('source_short', 'N/A')}"
        p_m.font.size, p_m.font.color.rgb = Pt(11), RGBColor(120, 120, 120)

    # 4. 保存 PPT
    output_name = f"Card_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(output_name)
    print(f"Success: [PPT卡片] 与 [Obsidian原子卡片] 已同步生成！")
    print(f"Obsidian 画布已更新: Inbox_Dashboard.canvas")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ppt_exporter.py <json_path>")
    else:
        create_ppt(sys.argv[1])